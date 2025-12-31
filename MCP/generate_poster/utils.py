from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from PIL import Image as PILImage
from pathlib import Path
import logging
logger = logging.getLogger(__name__)


def create_vertical_poster(text_summary: dict[str, str], image: PILImage) -> None:
    """Create a vertical academic-style poster."""
    
    # Create presentation with custom slide size (vertical poster)
    prs = Presentation()
    
    # Set slide dimensions: 24" wide x 36" tall (typical poster size)
    prs.slide_width = Inches(24)
    prs.slide_height = Inches(36)
    
    # Add a blank slide
    blank_slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(blank_slide_layout)
    
    # Define colors
    TITLE_COLOR = RGBColor(26, 35, 126)  # Dark blue
    ACCENT_COLOR = RGBColor(74, 144, 226)  # Light blue
    TEXT_COLOR = RGBColor(50, 50, 50)  # Dark gray
    
    # 1. ADD TITLE
    title_box = slide.shapes.add_textbox(
        left=Inches(1),
        top=Inches(0.5),
        width=Inches(22),
        height=Inches(2)
    )
    
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    title_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    
    title_p = title_frame.paragraphs[0]
    title_p.text = text_summary["title"]
    title_p.alignment = PP_ALIGN.CENTER
    title_p.font.size = Pt(72)
    title_p.font.bold = True
    title_p.font.color.rgb = TITLE_COLOR
    
    # 2. ADD DECORATIVE LINE UNDER TITLE
    line = slide.shapes.add_shape(
        1,  # Line shape
        left=Inches(2),
        top=Inches(2.8),
        width=Inches(20),
        height=Inches(0.1)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT_COLOR
    line.line.fill.background()
    
    # 3. ADD TEXT BOX WITH SECTIONS
    text_box = slide.shapes.add_textbox(
        left=Inches(1),
        top=Inches(3.5),
        width=Inches(22),
        height=Inches(18)
    )
    
    text_frame = text_box.text_frame
    text_frame.word_wrap = True
    text_frame.margin_left = Inches(0.5)
    text_frame.margin_right = Inches(0.5)
    text_frame.margin_top = Inches(0.3)
    
    # Section: Author
    author_p = text_frame.paragraphs[0]
    author_p.text = "Authors"
    author_p.font.size = Pt(48)
    author_p.font.bold = True
    author_p.font.color.rgb = ACCENT_COLOR
    author_p.space_after = Pt(12)
    
    # Author details
    author_detail = text_frame.add_paragraph()
    author_detail.text = ", ".join(text_summary["authors"])
    author_detail.font.size = Pt(36)
    author_detail.font.color.rgb = TEXT_COLOR
    author_detail.space_after = Pt(8)
    
    # Affiliations
    affiliation = text_frame.add_paragraph()
    affiliation.text = ", ".join(set(text_summary["affilications"]))
    affiliation.font.size = Pt(30)
    affiliation.font.italic = True
    affiliation.font.color.rgb = RGBColor(100, 100, 100)
    affiliation.space_after = Pt(40)
    
    # Section: Main Bullet Points
    section_b = text_frame.add_paragraph()
    section_b.text = "Key Findings"
    section_b.font.size = Pt(48)
    section_b.font.bold = True
    section_b.font.color.rgb = ACCENT_COLOR
    section_b.space_after = Pt(20)
    
    # Bullet points
    bullets = text_summary["points"]
    
    for bullet_text in bullets:
        bullet = text_frame.add_paragraph()
        bullet.text = bullet_text
        bullet.level = 0
        bullet.font.size = Pt(34)
        bullet.font.color.rgb = TEXT_COLOR
        bullet.space_after = Pt(16)
        
        # Add bullet symbol
        bullet.bullet = True
    
    # Add some spacing
    spacer = text_frame.add_paragraph()
    spacer.text = ""
    spacer.space_after = Pt(30)
    
    # Section: Image
    image.save("tmp.png")
    slide.shapes.add_picture(
        "tmp.png",
        left=Inches(6),
        top=Inches(16),
        width=Inches(12),
        height=Inches(16)
    )
    
    # Save the presentation
    if not Path("output").exists():
        Path("output").mkdir()
    prs.save(f'output/{text_summary["title"]}.pptx')
    logger.info(f"Poster created successfully: {text_summary["title"]}.pptx")
